from fastapi import APIRouter, UploadFile, File, HTTPException
from typing import List, Dict, Any
import io
import pypdf
from pptx import Presentation

router = APIRouter(prefix="/api/v1/documents", tags=["Documents"])

@router.post("/upload")
async def upload_document(file: UploadFile = File(...)):
    if not file.filename:
        raise HTTPException(status_code=400, detail="No file uploaded")
    
    file_ext = file.filename.split(".")[-1].lower()
    content = await file.read()
    
    extracted_data = []
    
    try:
        if file_ext == "txt":
            text = content.decode("utf-8", errors="ignore")
            extracted_data.append({
                "page_number": 1,
                "text": text,
                "metadata": {"type": "txt"}
            })
            
        elif file_ext == "pdf":
            pdf_file = io.BytesIO(content)
            reader = pypdf.PdfReader(pdf_file)
            for i, page in enumerate(reader.pages):
                text = page.extract_text()
                if text:
                    extracted_data.append({
                        "page_number": i + 1,
                        "text": text,
                        "metadata": {"type": "pdf_page"}
                    })
                    
        elif file_ext == "pptx":
            pptx_file = io.BytesIO(content)
            presentation = Presentation(pptx_file)
            for i, slide in enumerate(presentation.slides):
                text_runs = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
                text = "\n".join(text_runs)
                if text.strip():
                    extracted_data.append({
                        "page_number": i + 1,
                        "text": text,
                        "metadata": {"type": "pptx_slide"}
                    })
        else:
            raise HTTPException(status_code=400, detail=f"Unsupported file type: {file_ext}")
            
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing file: {str(e)}")
        
    return {
        "filename": file.filename,
        "document_type": file_ext,
        "pages": extracted_data
    }

from backend.app.services.chunking import DocumentChunker

@router.post("/chunks")
async def extract_and_chunk_document(file: UploadFile = File(...), chunk_size: int = 1000, chunk_overlap: int = 150):
    if not file.filename:
        raise HTTPException(status_code=400, detail="No file uploaded")
    
    file_ext = file.filename.split(".")[-1].lower()
    content = await file.read()
    
    extracted_data = []
    
    try:
        if file_ext == "txt":
            text = content.decode("utf-8", errors="ignore")
            extracted_data.append({
                "page_number": 1,
                "text": text,
                "metadata": {"type": "txt"}
            })
            
        elif file_ext == "pdf":
            pdf_file = io.BytesIO(content)
            reader = pypdf.PdfReader(pdf_file)
            for i, page in enumerate(reader.pages):
                text = page.extract_text()
                if text:
                    extracted_data.append({
                        "page_number": i + 1,
                        "text": text,
                        "metadata": {"type": "pdf_page"}
                    })
                    
        elif file_ext == "pptx":
            pptx_file = io.BytesIO(content)
            presentation = Presentation(pptx_file)
            for i, slide in enumerate(presentation.slides):
                text_runs = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
                text = "\n".join(text_runs)
                if text.strip():
                    extracted_data.append({
                        "page_number": i + 1,
                        "text": text,
                        "metadata": {"type": "pptx_slide"}
                    })
        else:
            raise HTTPException(status_code=400, detail=f"Unsupported file type: {file_ext}")
            
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing file: {str(e)}")
        
    chunker = DocumentChunker(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    chunks = chunker.chunk_document(file.filename, file_ext, extracted_data)
    
    # Return a summary and the chunks. We can limit the chunks returned if they're too large,
    # but for in-memory milestone, we return all or at least a structure containing them.
    return {
        "filename": file.filename,
        "total_extracted_sections": len(extracted_data),
        "total_chunks": len(chunks),
        "chunks": chunks
    }

from backend.app.services.embeddings import LocalEmbeddingService

embedding_service = LocalEmbeddingService()

@router.post("/embeddings")
async def extract_chunk_and_embed_document(file: UploadFile = File(...), chunk_size: int = 1000, chunk_overlap: int = 150):
    if not file.filename:
        raise HTTPException(status_code=400, detail="No file uploaded")
    
    file_ext = file.filename.split(".")[-1].lower()
    content = await file.read()
    
    extracted_data = []
    
    try:
        if file_ext == "txt":
            text = content.decode("utf-8", errors="ignore")
            extracted_data.append({
                "page_number": 1,
                "text": text,
                "metadata": {"type": "txt"}
            })
            
        elif file_ext == "pdf":
            pdf_file = io.BytesIO(content)
            reader = pypdf.PdfReader(pdf_file)
            for i, page in enumerate(reader.pages):
                text = page.extract_text()
                if text:
                    extracted_data.append({
                        "page_number": i + 1,
                        "text": text,
                        "metadata": {"type": "pdf_page"}
                    })
                    
        elif file_ext == "pptx":
            pptx_file = io.BytesIO(content)
            presentation = Presentation(pptx_file)
            for i, slide in enumerate(presentation.slides):
                text_runs = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
                text = "\n".join(text_runs)
                if text.strip():
                    extracted_data.append({
                        "page_number": i + 1,
                        "text": text,
                        "metadata": {"type": "pptx_slide"}
                    })
        else:
            raise HTTPException(status_code=400, detail=f"Unsupported file type: {file_ext}")
            
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing file: {str(e)}")
        
    chunker = DocumentChunker(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    chunks = chunker.chunk_document(file.filename, file_ext, extracted_data)
    
    embedded_chunks = embedding_service.embed_chunks(chunks)
    
    # Return a summary and the embedded chunks
    return {
        "filename": file.filename,
        "total_chunks": len(embedded_chunks),
        "embedding_dimension": embedding_service.embedding_dimension,
        "preview": embedded_chunks[:2] if embedded_chunks else []
    }
